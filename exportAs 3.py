'''
Author: Lalit Gupta
Date: 2024-03-22
Description: This script exports the Tableau workbook to PDF or PPT based on the configuration provided in the JSON file.
The script performs the following steps:
1. Connect to the Tableau Server.
2. Load the configuration from the specified file.
3. Retrieve the workbook ID based on the given workbook name and project name.
4. Apply the filters to the image export.
5. Export the images for each view in the workbook.
6. Export the workbook to PDF or PPT based on the specified export type.
7. Save the exported file to the specified location.
'''

import tableauserverclient as TSC
import datetime
import requests
import os
import argparse
import json
import uuid
from pptx import Presentation
from pptx.util import Inches
from PIL import Image
from fpdf import FPDF

# update for new requirements
from PIL import Image

# Configurtion
server_url = 'https://prod-useast-b.online.tableau.com/'
site_name = 'b360bi'
token_name = 'myToken'
access_token = 'PsfRRSmjSJuLs8Nn5BhnBA==:bE8wajmFG2KOslgT1FgRnQKdVyBEcKoC'
api_version = '3.22'

# ppt_location = r'C:\Workspace\Automations\content'
ppt_location = ''
img_export = "img/"
img_indexing = []


def connect_tableau():
    # Step 2: Connect to the Tableau Server
    try:
        tableau_auth = TSC.PersonalAccessTokenAuth(token_name=token_name,
                                                   personal_access_token=access_token,
                                                   site_id=site_name)
        server = TSC.Server(server_url)
        server.auth.sign_in(tableau_auth)
        print("DEBUG: Connected to Tableau Server")
        return server
    except Exception as err:
        print("ERR: Error connecting to Tableau Server")
        raise err

def get_workbook(server, workbook_name, project_name):
    """ Retrieves the workbook ID based on the given workbook name and project name.

    Args:
        server (TableauServer): The Tableau Server object.
        workbook_name (str): The name of the workbook to search for.
        project_name (str): The name of the project where the workbook resides.

    Returns:
        int: The ID of the matching workbook.

    Raises:
        FileExistsError: If no workbook with the specified name is found or if multiple workbooks
                        share the same name.
    """
    # Get all workbooks in the specified project
    all_workbooks = list(TSC.Pager(server.workbooks))
    matching_workbooks = [workbook for workbook in all_workbooks if
                          workbook.name == workbook_name and workbook.project_name == project_name]
    if len(matching_workbooks) == 0:
        error = f"Workbook '{workbook_name}' not found in project '{project_name}'."
        raise FileExistsError(error)
    elif len(matching_workbooks) > 1:
        error = f"Multiple workbooks with name '{workbook_name}' found in project '{project_name}'."
        raise FileExistsError(error)

    return matching_workbooks[0]


def load_config(config_file):
    """ Loads the configuration from the specified file.

    Args:
        config_file (str): The path to the configuration file.

    Returns:
        dict: The configuration settings.
    """
    try:
        with open(config_file, "r") as file:
            config = json.load(file)
            return config
    except Exception as err:
        print(f"ERR: Error loading configuration file '{config_file}'.")
        raise err


def get_image_export(filters) -> TSC.ImageRequestOptions: 
    """ Get the image export object with filters.

    Args:
        filters (dict): The filters to apply to the views.

    Returns:
        TSC.ImageRequestOptions: The image export object with filters applied.
    """
    print("DEBUG: Applying filters to the image export.")
    image_export = TSC.ImageRequestOptions(imageresolution=TSC.ImageRequestOptions.Resolution.High, maxage=1)
    print("DEBUG: Filters Details: ")
    if filters:
        for filter in filters:
            filter_disc = filter['description']
            filter_values = ','.join(filter['values'])
            print(f"-filter '{filter_disc}' with values '{filter_values}'")
            image_export.vf(filter_disc, filter_values)
    else:
        print("DEBUG: No filters to apply.")
            
    return image_export


def export_images(server, workbook, image_export_option, img_dir, ppt_slides):
    """ Export the images for each view in the workbook.

    Args:
        server (TableauServer): The Tableau Server object.
        workbook (TableauWorkbook): The workbook to export.
        image_export_option (TSC.ImageRequestOptions): The image export options.
        img_dir (str): The directory to save the images to.
    """
    for view in workbook.views:
        view_name = view.name
        # Skip the views that are not in the list of slides to export
        if view_name not in ppt_slides:
            continue
        print(f"DEBUG: Exporting image for view '{view_name}'")
        server.views.populate_image(view, image_export_option)
        image = view.image
        image_path = os.path.join(img_dir, f"{view_name}.png")
        img_indexing.append(f"{view_name}.png")
        with open(image_path, "wb") as file:
            file.write(image)
        print(f"DEBUG: Image exported to '{image_path}'")

        server.views.populate_csv(view, image_export_option)
        #csv_path = file_path1 + view.name + '.csv'
        csv_path = os.path.join(img_dir, f"{view_name}.csv")
        with open(csv_path, 'wb') as csv_file:
            # Perform byte join on the CSV data
            csv_file.write(b''.join(view.csv))


def get_images(img_dir):
    """ Get the list of exported images.

    Args:
        img_dir (str): The directory containing the exported images.

    Returns:
        list: The list of image files in the directory.
    """
    # img_files = [file for file in os.listdir(img_dir) if file.endswith(".png")]
    # return img_files
    return img_indexing


def add_title_slide(prs, workbook_name, project_name):
    """ Add a title slide to the PowerPoint presentation.

    Args:
        prs (Presentation): The PowerPoint presentation object.
        workbook_name (str): The name of the workbook being exported.
        export_path (str): The path where the exported file will be saved.
        img_dir (str): The directory containing the exported images.
    """
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = f"{workbook_name}"
    # Add project name and export date as subtitle
    subtitle.text = f"Project: {project_name}\nExported on: {datetime.datetime.now()}"
    

def add_image_to_ppt(prs, img_dir, img_file):
    """ Add an image to the PowerPoint presentation.

    Args:
        prs (Presentation): The PowerPoint presentation object.
        img_dir (str): The directory containing the exported images.
        img_file (str): The name of the image file to add.
    """
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.placeholders[0]
    tp = title.element
    tp.getparent().remove(tp)

    # Get the dimensions of the slide
    slide_width, slide_height = prs.slide_width, prs.slide_height

    # Load the image
    img_path = os.path.join(img_dir, img_file)
    img = Image.open(img_path)
    img_width, img_height = img.size

    # Calculate the scaling factor to fit the image within the slide
    max_width = Inches(8)
    max_height = Inches(6)
    width_scale = max_width / img_width
    height_scale = max_height / img_height
    scale = min(width_scale, height_scale)

    # Calculate the centered position
    left = (slide_width - img_width * scale) / 2
    top = (slide_height - img_height * scale) / 2

    # Add the image to the slide
    slide.shapes.add_picture(img_path, left, top, width=img_width * scale, height=img_height * scale)
    print(f"DEBUG: Added image '{img_file}' to the PowerPoint presentation")


def create_deck(img_dir, export_path, workbook_name):
    """ Create a PPT deck from the exported images.

    Args:
        img_dir (str): The directory containing the exported images.
        export_path (str): The path to save the PPT deck to.
        workbook_name (str): The name of the workbook being exported.
    """

    # Create a new PowerPoint presentation
    base_ppt  = "base.pptx"   # TODO: Use a base template for the PPT deck
    prs = Presentation()
    print(f"DEBUG: Created new PowerPoint presentation")

    # Add a title slide to the presentation
    add_title_slide(prs, workbook_name, project_name)
    

    # Add the exported images to the presentation
    img_files = get_images(img_dir)
    print(f"DEBUG: Adding images to the PowerPoint presentation")
    # print(img_files)
    for img_file in img_files:
        add_image_to_ppt(prs, img_dir, img_file)

    prs.save(export_path)
    print(f"DEBUG: Saved PowerPoint presentation to '{export_path}'")


def add_title_page(pdf, pdf_config):
    """ Add a title page to the PDF deck.

    Args:
        pdf (FPDF): The PDF object to add the title page to.
        pdf_config (dict): The configuration settings for the PDF export.
    """
    # Step 1: Add a new page and set the config for the title
    pdf.add_page()
    pdf.set_font("Arial", style="B", size=24)

    # Calculate the y-coordinate to place the text at the center of the page
    page_height = pdf.h
    cell_height = pdf.font_size    
    y = (page_height - cell_height) / 2

    title = f"Tableau Workbook Export: {pdf_config['workbook_name']}"
    # Calculate the x-coordinate to place the text at the center of the page
    page_width = pdf.w
    cell_width = pdf.get_string_width(title) + 2 * pdf.c_margin
    x = (page_width - cell_width) / 2
    pdf.set_xy(x, y)
    pdf.cell(cell_width, 10, title, ln=True, align="C")

    # Step 2: Add the project name and export date
    pdf.set_font("Arial", style="I", size=16)
    if "project_name" in pdf_config:
        project_subtitle = f"Project: {pdf_config['project_name']}"
        cell_width = pdf.get_string_width(project_subtitle) + 2 * pdf.c_margin
        x = (page_width - cell_width) / 2
        pdf.set_x(x)
        pdf.cell(cell_width, 10, project_subtitle, ln=True, align="C")
        # pdf.cell(200, 10, f"Project: {pdf_config['project_name']}", ln=True, align="C")

    pdf.set_font("Arial", style="I", size=12)
    datetime_info = f"Exported on: {datetime.datetime.now()}"
    cell_width = pdf.get_string_width(datetime_info) + 2 * pdf.c_margin
    x = (page_width - cell_width) / 2
    pdf.set_x(x)
    pdf.cell(cell_width, 10, datetime_info, ln=True, align="C")
    # pdf.cell(200, 10, f"Exported on: {datetime.datetime.now()}", ln=True, align="C")


def add_image_to_pdf(pdf, img_path, pdf_config):
    """ Add an image to the PDF deck.

    Args:
        pdf (FPDF): The PDF object to add the image to.
        img_path (str): The path to the image file.
        pdf_config (dict): The configuration settings for the PDF export.
    """
    img = Image.open(img_path)
    img_width, img_height = img.size
    # Calculate the image dimensions to fit within the page
    page_width, page_height = pdf.w, pdf.h
    if img_width > page_width or img_height > page_height:
        ratio_width = page_width / img_width
        ratio_height = page_height / img_height
        if ratio_width < ratio_height:
            new_width = page_width
            new_height = img_height * ratio_width
        else:
            new_width = img_width * ratio_height
            new_height = page_height
    else:
        new_width = img_width
        new_height = img_height

    # Calculate the x and y coordinates to place the image at the center of the page
    x = (page_width - new_width) / 2
    y = (page_height - new_height) / 2
    
    pdf.add_page()
    pdf.image(img_path, x=x, y=y, w=new_width, h=new_height)


def export_as_pdf(img_dir, output_path, pdf_config):
    """ Create a PDF deck from the exported images.

    Args:
        img_dir (str): The directory containing the exported images.
        output_path (str): The path to save the PPT deck to.
    """

    # Step 1: Create a new PDF object with the specified configuration
    print("DEBUG: Creating PDF deck")
    pdf = FPDF(orientation=pdf_config["orientation"], unit="mm", format=pdf_config["page_type"])
    pdf.set_auto_page_break(auto=True, margin=15)
    # Step 2: Add the images to the PDF deck
    add_title_page(pdf, pdf_config)
    # Step 3.1: Get the list of exported images  
    img_files = get_images(img_dir)
    # Step 3.2: Add each image to the PDF deck
    for img_file in img_files:
        img_path = os.path.join(img_dir, img_file)
        add_image_to_pdf(pdf, img_path, pdf_config)      
    # Step 4: Save the PDF deck to the specified location
    pdf.output(output_path)    


def split_image(image_path, output_dir, rows, cols, img_name, tmp_index):
    """
    Split the image into tiles and save them to the output directory

    Args:
        image_path (str): The path to the image file to split.
        output_dir (str): The directory to save the tiles to.
        rows (int): The number of rows to split the image into.
        cols (int): The number of columns to split the image into.
    """
    img = Image.open(image_path)
    width, height = img.size
    
    tile_width = width // cols
    tile_height = height // rows
    
    os.makedirs(output_dir, exist_ok=True)
    for row in range(rows):
        for col in range(cols):
            if row==0 and col==1: # Skip the second tile in the first row
                continue
            if row==0:
                left = 0
                upper = row * tile_height
                right = width
                lower = upper + tile_height
            else:
                left = col * tile_width
                upper = row * tile_height
                right = left + tile_width
                lower = upper + tile_height
                
            tile = img.crop((left, upper, right, lower))
            export_name = f"{img_name}_tile_{row}_{col}.png" 
            tile_path = os.path.join(output_dir, export_name)
            tile.save(tile_path)
            tmp_index.append(export_name)
            # print(f"Saved {tile_path}")


def generate_slide_imgs(img_dir):
    """ Split the images into tiles for respective slides.

    Args:
        img_dir (str): The directory containing the exported images.
        img_indexing (list): The list of image files in the directory.
    """
    global img_indexing
    tmp_index = []
    for img_file in img_indexing:
        img_path = os.path.join(img_dir, img_file)
        export_dir = os.path.join(img_dir, "slide_imgs")
        img_name = img_file.split(".")[0]
        split_image(img_path, export_dir, 2, 2, img_name, tmp_index)             
    print(f"DEBUG: Split images into tiles for respective slides")
    img_indexing = tmp_index # Update the image indexing list
    return os.path.join(img_dir, "slide_imgs") # Update the image directory
    # print(f"DEBUG: Updated image indexing list: {img_indexing}")


if __name__ == "__main__":

    # STEP 0: Collect the command-line arguments
    parser = argparse.ArgumentParser(description="Export a Tableau workbook to PDF or PPT.")
    parser.add_argument("config_file", help="The path to the configuration file.")
    parser.add_argument("export_type", choices=["pdf", "ppt"], help="The type of export to perform.")
    # Config for pdf  e.g. page_type, orientation
    parser.add_argument("--page_type", default="A4", help="The page type for the PDF (e.g., 'A4').")
    parser.add_argument("--orientation", default="landscape", help="The orientation for the PDF (e.g., 'landscape').")
    
    args = parser.parse_args()
    config_file = "config/"+args.config_file
    export_type = args.export_type
    # PDF specific arguments
    # Page type can be: A3, A4, A5, B5, Executive, Folio, Ledger, Legal, Letter, Note, Quarto, or Tabloid.
    # Orientation can be: landscape or portrait
    page_type = args.page_type
    orientation = args.orientation
    pdf_config = {"page_type": page_type, "orientation": orientation}
    
    # STEP 1: Connect to Tableau Server
    server = connect_tableau()
    server.version = '3.22'

    # STEP 2: Process the configuration for the export
    config = load_config(config_file)
    workbook_name = config["workbook_name"]
    project_name = config["project_name"]
    ppt_slides = config["slide_views"]


    print("----------------------------------------STEP#1------------------------------------------------")
    print(f"DEBUG: Configuration loaded:\n- Workbook: {workbook_name}\n- Project: {project_name}")
    # STEP 3: Get the workbook luid to be exported
    workbook = get_workbook(server, workbook_name, project_name)
    # STEP 3.1: Populate the views in the workbook
    server.workbooks.populate_views(workbook)
    # print("DEBUG: Workbook views populated")

    # STEP 4: SetUP Image export with filters
    print("----------------------------------------------------------------------------------------")
    image_export_option =  get_image_export(config['filters'])
    print("----------------------------------------------------------------------------------------")
    # STEP 5: Export image into the img_dir for each view
    # STEP 5.1: Create the directory with a unique id
    export_id = str(uuid.uuid4())
    img_dir = os.path.join(img_export, export_id)
    os.makedirs(img_dir, exist_ok=True)
    print(f"DEBUG: Image directory created: {img_dir}")
    
    # STEP 5.2: Export the images for each view
    export_images(server, workbook, image_export_option, img_dir, ppt_slides)
    # STEP 5.3: Split the images for respective slides
    img_dir = generate_slide_imgs(img_dir)
    print("----------------------------------------STEP#2------------------------------------------------")

    # quit()
    if export_type == "ppt":
        # STEP 6: Export the workbook to PPT
        print("DEBUG: Exporting workbook to PPT")
        ppt_location = os.path.join(ppt_location,'ppt')
        os.makedirs(ppt_location, exist_ok=True)
        export_path = os.path.join(ppt_location, f"{workbook_name}.pptx")
        print(f"DEBUG: Exporting workbook to '{export_path}'")
        create_deck(img_dir, export_path, workbook_name)
        
    if export_type == "pdf":
        print("DEBUG: Exporting workbook to PDF")
        # Add workbook name and project name to the pdf config
        pdf_config["workbook_name"] = workbook_name
        pdf_config["project_name"] = project_name
        pdf_dir = os.path.join(ppt_location, "pdf")  # Update the directory name
        os.makedirs(pdf_dir, exist_ok=True)  # Ensure the directory exists or create it
        output_file = os.path.join(pdf_dir, f"{workbook_name}.pdf")
        export_as_pdf(img_dir, output_file, pdf_config)
        print(f"DEBUG: Exported workbook '{workbook_name}' as PDF to '{output_file}'")


    # STEP 7: Clean up the image directory
    print("DEBUG: Cleaning up the image directory")
    for file in os.listdir(img_dir):
        file_path = os.path.join(img_dir, file)
        os.remove(file_path)
    os.rmdir(img_dir)
    print("DEBUG: Image directory cleaned up.")

    # STEP 8: Sign out from the Tableau Server
    print("DEBUG: Signing out from Tableau Server")  
    server.auth.sign_out()  
    

